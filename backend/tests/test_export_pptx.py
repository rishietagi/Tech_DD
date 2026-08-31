"""PowerPoint export.

Unlike the PDF, a .pptx is a zip of XML that python-pptx can read straight back, so
these assert on the real object model — slide count, shape text, table contents —
rather than on scraped bytes.
"""

import io

from fastapi.testclient import TestClient
from pptx import Presentation

from app.services.scope.export_pptx import (
    build_deck,
    pptx_filename,
    render_pptx,
    render_pptx_from_payload,
)
from app.services.scope.rules_generator import RulesScopeGenerator
from tests.factories import make_intake
from tests.test_engagement_flow import SECTIONS_PAYLOAD

PRODUCT = dict(
    company_name="Meridian Analytics",
    line_of_business="Sells a usage-based analytics platform to mid-market e-commerce retailers.",
    dd_type_preference="Product Tech DD",
    digital_maturity="Digital native",
    data_sensitivity=["Personal data (PII)"],
    compliance_regimes=["SOC 2"],
)

PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def scope_deck(**overrides):
    scope = RulesScopeGenerator().generate(make_intake(**{**PRODUCT, **overrides}))
    return build_deck(scope, deal_name="Project Lighthouse", version=1)


def deck_text(prs) -> str:
    """Every run of text in the deck, flattened. Includes table cells."""
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    chunks.extend(cell.text for cell in row.cells)
    return "\n".join(chunks)


def _filed_engagement(client: TestClient) -> str:
    response = client.post("/api/v1/engagements", json={"deal_name": "Project Lighthouse"})
    engagement_id = response.json()["id"]
    for section, payload in SECTIONS_PAYLOAD.items():
        client.patch(f"/api/v1/engagements/{engagement_id}/intake/{section}", json=payload)
    client.post(f"/api/v1/engagements/{engagement_id}/submit")
    return engagement_id


# ------------------------------------------------------------------------ renderer


def test_render_produces_a_readable_pptx() -> None:
    raw = render_pptx(
        RulesScopeGenerator().generate(make_intake(**PRODUCT)), "Project Lighthouse", 1
    )
    assert raw.startswith(b"PK")  # a .pptx is a zip
    reopened = Presentation(io.BytesIO(raw))
    assert len(reopened.slides) >= 5  # cover, engagement, rows, cost, team/exclusions


def test_deck_is_16_by_9() -> None:
    prs = scope_deck()
    assert round(prs.slide_width / prs.slide_height, 2) == round(16 / 9, 2)


def test_cover_carries_the_deal_and_the_classification() -> None:
    text = deck_text(scope_deck())
    assert "Project Lighthouse" in text
    assert "Product Due Diligence" in text
    assert "SCOPE OF WORK" in text
    assert "Classification" in text


def test_every_scope_row_appears() -> None:
    scope = RulesScopeGenerator().generate(make_intake(**PRODUCT))
    text = deck_text(build_deck(scope, "Project Lighthouse", 1))
    for row in scope.rows:
        assert row.title in text, f"row missing from the deck: {row.id}"


def test_client_facing_sections_are_present() -> None:
    text = deck_text(scope_deck())
    for heading in ("Engagement", "Scope of work", "Sequencing", "Cost estimation", "Exclusions"):
        assert heading in text, f"missing section: {heading}"


def test_the_internal_audit_trail_is_omitted() -> None:
    """Same split as the PDF: the deck is client-facing, Markdown is internal."""
    text = deck_text(scope_deck())
    assert "Provenance" not in text
    assert "Signals" not in text


def test_sequencing_shows_the_broad_pass_handoff() -> None:
    """The deck must carry the output of each phase, not just its name."""
    text = deck_text(scope_deck())
    assert "Broad pass" in text
    assert "areas of focus" in text


def test_exclusions_are_never_dropped() -> None:
    assert "does not say what it excludes" in deck_text(scope_deck())


def test_special_characters_survive() -> None:
    """python-pptx escapes XML itself, but a company name with & must still round-trip."""
    scope = RulesScopeGenerator().generate(
        make_intake(**{**PRODUCT, "company_name": "Ampersand & Sons <Holdings>"})
    )
    text = deck_text(build_deck(scope, "Deal & Co <Alpha>", 1))
    assert "Deal & Co <Alpha>" in text


def test_only_schema_v2_is_exportable() -> None:
    try:
        render_pptx_from_payload({"schema_version": 1}, "Project Lighthouse", 1)
    except ValueError as exc:
        assert "schema v2" in str(exc)
    else:
        raise AssertionError("a v1 payload must not be exportable")


def test_filename_is_slugged() -> None:
    assert pptx_filename("Project Redline", 2) == "project-redline-scope-v2.pptx"
    assert pptx_filename("!!!", 1) == "scope-v1.pptx"


# -------------------------------------------------------------------------- route


def test_export_pptx_endpoint_returns_a_deck(client: TestClient) -> None:
    engagement_id = _filed_engagement(client)
    client.post(f"/api/v1/engagements/{engagement_id}/scope")

    response = client.get(f"/api/v1/engagements/{engagement_id}/scope/export.pptx")
    assert response.status_code == 200
    assert response.headers["content-type"] == PPTX_MEDIA_TYPE
    assert response.content.startswith(b"PK")

    disposition = response.headers["content-disposition"]
    assert "attachment" in disposition
    assert "project-redline-scope-v1.pptx" in disposition


def test_export_pptx_route_is_not_shadowed_by_the_version_route(client: TestClient) -> None:
    """`/export.pptx` is declared before `/{version}`; confirm it is not read as one."""
    engagement_id = _filed_engagement(client)
    client.post(f"/api/v1/engagements/{engagement_id}/scope")

    response = client.get(f"/api/v1/engagements/{engagement_id}/scope/export.pptx")
    assert response.status_code == 200
    assert response.headers["content-type"] == PPTX_MEDIA_TYPE


def test_export_pptx_without_a_scope_returns_404(client: TestClient) -> None:
    engagement_id = _filed_engagement(client)
    response = client.get(f"/api/v1/engagements/{engagement_id}/scope/export.pptx")
    assert response.status_code == 404
    assert response.json()["detail"]["code"] == "no_scope"


def test_human_edits_survive_into_the_deck(client: TestClient) -> None:
    engagement_id = _filed_engagement(client)
    scope = client.post(f"/api/v1/engagements/{engagement_id}/scope").json()
    row_id = scope["payload"]["rows"][0]["id"]

    client.patch(
        f"/api/v1/engagements/{engagement_id}/scope/1/rows/{row_id}",
        json={"tier": 1, "reason": "Client has had this area reviewed externally."},
    )

    response = client.get(f"/api/v1/engagements/{engagement_id}/scope/export.pptx")
    text = deck_text(Presentation(io.BytesIO(response.content)))
    # The edited row still ships; only its depth changed, which the deck does not label.
    edited = next(r for r in scope["payload"]["rows"] if r["id"] == row_id)
    assert edited["title"] in text


def test_tier_badges_are_not_shown_on_the_scope_slides() -> None:
    """Depth is an internal decision (Rishi, 2026-08-31) and is not badged on rows.

    Scoped to the scope-of-work slides only: "Broad pass" and "Deep dive" are the
    sequencing *phase names* and must still appear on the sequencing slide.
    """
    prs = scope_deck()
    row_slides = [
        s for s in prs.slides
        if any(
            sh.has_text_frame and sh.text_frame.text.startswith("Scope of work")
            for sh in s.shapes
        )
    ]
    assert row_slides, "no scope-of-work slides found"

    for slide in row_slides:
        text = " ".join(
            sh.text_frame.text for sh in slide.shapes if sh.has_text_frame
        ).upper()
        for label in ("DEEP DIVE", "ASSESS", "SCREEN", "SWEEP"):
            assert label not in text, f"tier badge leaked onto a scope slide: {label}"
