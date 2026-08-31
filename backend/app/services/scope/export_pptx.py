"""PowerPoint export of a generated scope — the client-facing deck.

A scope of work is usually circulated as a deck rather than a document, so this is the
primary client-facing artefact; `export_pdf.py` is the same content laid out as a paper
document, and `export.py` (Markdown) is the internal one that carries the audit trail.

Content matches the PDF exactly: the internal reasoning layer (signals, rule
provenance) is omitted; exclusions are not internal and are always present.

**Styling is deliberately separated from content.** Every colour, size and font lives in
`Theme`; `build_deck` only decides what goes on which slide. When a KPMG template .pptx
arrives, the change is to load it as the presentation base and map `Theme` onto its
masters — the slide-building code below does not need to change.

Deterministic and offline: read off the stored payload, no recomputation, no network.
"""

import io
import re
from dataclasses import dataclass
from datetime import date
from pathlib import Path
from typing import Any

# The package-level `Presentation` is a factory function, not a class, so the real
# class is imported separately for annotations.
from pptx import Presentation as open_presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation
from pptx.util import Emu, Inches, Pt

from app.schemas.scope import ScopeOfWorkPayloadV2

_DECK_LABEL = {
    "product": "Product Tech DD",
    "enterprise": "Enterprise IT DD",
    "blended": "Blended",
}

_LOGO = Path(__file__).resolve().parents[4] / "assets" / "kpmg-logo-blue.png"


@dataclass(frozen=True)
class Theme:
    """Every visual decision in one place, so a template swap is a contained change.

    Colours mirror `frontend/src/styles/tokens.css` and `export_pdf.py`, so the deck,
    the PDF and the screen are recognisably the same document.
    """

    blue: RGBColor = RGBColor(0x00, 0x33, 0x8D)
    blue_dark: RGBColor = RGBColor(0x00, 0x25, 0x66)
    blue_light: RGBColor = RGBColor(0x00, 0x91, 0xDA)
    text: RGBColor = RGBColor(0x0A, 0x0E, 0x1A)
    muted: RGBColor = RGBColor(0x4A, 0x55, 0x68)
    muted_2: RGBColor = RGBColor(0x80, 0x87, 0xA0)
    line: RGBColor = RGBColor(0xDD, 0xE3, 0xED)
    paper_2: RGBColor = RGBColor(0xF4, 0xF6, 0xFA)
    white: RGBColor = RGBColor(0xFF, 0xFF, 0xFF)

    font: str = "Segoe UI"

    # 16:9. Inches, matching python-pptx's own units.
    slide_w: float = 13.333
    slide_h: float = 7.5
    margin: float = 0.62

    @property
    def content_w(self) -> float:
        return self.slide_w - 2 * self.margin


THEME = Theme()

# How many scope rows fit on one slide before it needs to continue.
_ROWS_PER_SLIDE = 4


def _blank(prs: Presentation) -> Any:
    """A slide with no placeholders — every element is positioned explicitly.

    Layout 6 is the blank layout in the default template. With a KPMG template this
    becomes the branded content layout instead.
    """
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    anchor: Any = MSO_ANCHOR.TOP,
) -> Any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    return frame


def _write(
    frame: Any,
    text: str,
    *,
    size: float,
    bold: bool = False,
    color: RGBColor | None = None,
    italic: bool = False,
    space_after: float = 0,
    first: bool = False,
    theme: Theme = THEME,
    bullet: bool = False,
) -> Any:
    """Append a paragraph. `first` reuses the frame's initial empty paragraph."""
    para = frame.paragraphs[0] if first else frame.add_paragraph()
    run = para.add_run()
    run.text = f"•  {text}" if bullet else text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = theme.font
    run.font.color.rgb = color if color is not None else theme.text
    para.space_after = Pt(space_after)
    return para


def _accent_bar(slide: Any, theme: Theme = THEME) -> None:
    """The blue rule across the top of every content slide."""
    from pptx.enum.shapes import MSO_SHAPE

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(theme.slide_w), Inches(0.075)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.blue
    bar.line.fill.background()
    bar.shadow.inherit = False


def _slide_header(
    slide: Any, title: str, deal_name: str | None, theme: Theme = THEME
) -> float:
    """Accent bar, section title and running deal name. Returns the content top edge."""
    _accent_bar(slide, theme)

    frame = _textbox(slide, theme.margin, 0.34, theme.content_w, 0.5)
    _write(frame, title, size=24, bold=True, color=theme.blue, first=True, theme=theme)

    if deal_name:
        rf = _textbox(slide, theme.margin, 0.30, theme.content_w, 0.28)
        para = _write(rf, deal_name, size=10, color=theme.muted_2, first=True, theme=theme)
        para.alignment = PP_ALIGN.RIGHT

    return 1.05


def _footer(slide: Any, text: str, page: int | None, theme: Theme = THEME) -> None:
    frame = _textbox(slide, theme.margin, theme.slide_h - 0.46, theme.content_w, 0.26)
    para = _write(frame, text, size=8, color=theme.muted_2, first=True, theme=theme)
    para.alignment = PP_ALIGN.LEFT

    if page is not None:
        pf = _textbox(slide, theme.slide_w - theme.margin - 0.6, theme.slide_h - 0.46, 0.6, 0.26)
        pp = _write(pf, str(page), size=8, color=theme.muted_2, first=True, theme=theme)
        pp.alignment = PP_ALIGN.RIGHT


def _table(
    slide: Any,
    rows: list[list[str]],
    col_widths: list[float],
    top: float,
    theme: Theme = THEME,
) -> Any:
    """A header-row table in house style. `rows[0]` is the header."""
    n_rows, n_cols = len(rows), len(col_widths)
    shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(theme.margin), Inches(top),
        Inches(sum(col_widths)), Inches(0.32 * n_rows),
    )
    table = shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    for r, row_values in enumerate(rows):
        for c, value in enumerate(row_values):
            cell = table.cell(r, c)
            cell.text = ""
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.TOP
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = theme.blue
            else:
                cell.fill.fore_color.rgb = theme.white if r % 2 else theme.paper_2

            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = value
            run.font.size = Pt(9 if r else 9.5)
            run.font.bold = r == 0
            run.font.name = theme.font
            run.font.color.rgb = theme.white if r == 0 else theme.text
    return table


# ------------------------------------------------------------------------ slides


def _cover_slide(
    prs: Presentation,
    scope: ScopeOfWorkPayloadV2,
    deal_name: str | None,
    version: int | None,
    theme: Theme = THEME,
) -> None:
    slide = _blank(prs)
    _accent_bar(slide, theme)

    if _LOGO.exists():
        slide.shapes.add_picture(
            str(_LOGO), Inches(theme.margin), Inches(0.75), width=Inches(1.5)
        )

    frame = _textbox(slide, theme.margin, 2.3, theme.content_w, 2.0)
    _write(frame, "SCOPE OF WORK", size=11, bold=True, color=theme.blue, first=True, theme=theme)
    _write(frame, scope.deck_title, size=40, bold=True, color=theme.blue, space_after=4, theme=theme)
    _write(frame, scope.deck_subtitle, size=16, color=theme.muted, theme=theme)

    if deal_name:
        df = _textbox(slide, theme.margin, 4.45, theme.content_w, 0.4)
        _write(df, deal_name, size=18, bold=True, first=True, theme=theme)

    c = scope.classification
    meta = [
        ("Classification", _DECK_LABEL.get(c.dd_type.value, c.dd_type.value)),
        ("Enterprise / Product mix", f"{c.dd_mix} / 100"),
        ("Confidence", c.confidence.capitalize()),
        ("Areas in scope", str(sum(1 for r in scope.rows if r.tier > 0))),
        ("Date", date.today().strftime("%d %B %Y")),
    ]
    if version is not None:
        meta.append(("Version", str(version)))

    mf = _textbox(slide, theme.margin, 5.05, theme.content_w, 1.8)
    for i, (label, value) in enumerate(meta):
        para = mf.paragraphs[0] if i == 0 else mf.add_paragraph()
        lr = para.add_run()
        lr.text = f"{label}:  "
        lr.font.size = Pt(10.5)
        lr.font.name = theme.font
        lr.font.color.rgb = theme.muted
        vr = para.add_run()
        vr.text = value
        vr.font.size = Pt(10.5)
        vr.font.bold = True
        vr.font.name = theme.font
        vr.font.color.rgb = theme.text
        para.space_after = Pt(3)


def _engagement_slide(
    prs: Presentation, scope: ScopeOfWorkPayloadV2, deal_name: str | None,
    footer: str, page: int, theme: Theme = THEME,
) -> None:
    slide = _blank(prs)
    top = _slide_header(slide, "Engagement", deal_name, theme)

    frame = _textbox(slide, theme.margin, top, theme.content_w, 1.5)
    _write(frame, scope.engagement_summary, size=12, color=theme.text, first=True, theme=theme)

    if scope.objectives:
        of = _textbox(slide, theme.margin, top + 1.65, theme.content_w, 4.0)
        _write(of, "Objectives", size=15, bold=True, color=theme.blue_dark, first=True,
               space_after=6, theme=theme)
        for objective in scope.objectives:
            _write(of, objective, size=11, color=theme.text, bullet=True, space_after=5, theme=theme)

    _footer(slide, footer, page, theme)


def _row_slides(
    prs: Presentation, scope: ScopeOfWorkPayloadV2, deal_name: str | None,
    footer: str, start_page: int, theme: Theme = THEME,
) -> int:
    """One slide per `_ROWS_PER_SLIDE` rows. Returns the next page number."""
    page = start_page
    rows = scope.rows
    chunks = [rows[i : i + _ROWS_PER_SLIDE] for i in range(0, len(rows), _ROWS_PER_SLIDE)]

    for index, chunk in enumerate(chunks):
        slide = _blank(prs)
        heading = "Scope of work" if index == 0 else f"Scope of work (cont. {index + 1})"
        top = _slide_header(slide, heading, deal_name, theme)

        row_h = (theme.slide_h - top - 0.7) / _ROWS_PER_SLIDE
        for i, row in enumerate(chunk):
            y = top + i * row_h

            # No tier badge (Rishi, 2026-08-31): depth is an internal decision and the
            # client-facing document states the work, not the engine's grading of it.
            # The title reclaims the width the chip used to occupy.
            tf = _textbox(slide, theme.margin, y, theme.content_w, 0.32)
            _write(tf, f"{row.sn:02d}.  {row.title}", size=12, bold=True, first=True, theme=theme)

            bf = _textbox(slide, theme.margin, y + 0.34, theme.content_w - 0.2, row_h - 0.42)
            first = True
            for line in row.lines:
                _write(bf, line.text, size=10, color=theme.muted, first=first,
                       space_after=3, theme=theme)
                first = False
            if row.out_of_scope_note:
                _write(bf, row.out_of_scope_note, size=9, color=theme.muted_2,
                       italic=True, first=first, theme=theme)

        _footer(slide, footer, page, theme)
        page += 1

    return page


def _sequencing_slide(
    prs: Presentation, scope: ScopeOfWorkPayloadV2, deal_name: str | None,
    footer: str, page: int, theme: Theme = THEME,
) -> None:
    slide = _blank(prs)
    top = _slide_header(slide, "Sequencing", deal_name, theme)

    sub = _textbox(slide, theme.margin, top - 0.06, theme.content_w, 0.3)
    _write(sub, "A broad pass identifies the areas of focus; the deep dive works on those.",
           size=10.5, color=theme.muted, italic=True, first=True, theme=theme)

    rows: list[list[str]] = [["Phase", "Weeks", "Focus", "Output"]]
    for phase in scope.sequencing:
        rows.append([phase.name, phase.weeks, phase.focus, phase.output or "—"])
    _table(slide, rows, [1.7, 1.1, 5.4, 3.9], top + 0.34, theme)

    _footer(slide, footer, page, theme)


def _cost_slide(
    prs: Presentation, scope: ScopeOfWorkPayloadV2, deal_name: str | None,
    footer: str, page: int, theme: Theme = THEME,
) -> None:
    slide = _blank(prs)
    top = _slide_header(slide, "Cost estimation", deal_name, theme)

    af = _textbox(slide, theme.margin, top, theme.content_w, 0.6)
    _write(af, scope.cost_plan.approach, size=11, color=theme.text, first=True, theme=theme)

    y = top + 0.75
    if scope.cost_plan.lines:
        rows: list[list[str]] = [["Type", "Line", "Basis"]]
        for line in scope.cost_plan.lines:
            rows.append([
                "One-time" if line.category == "one_time" else "Recurring",
                line.label,
                line.basis,
            ])
        _table(slide, rows, [1.4, 4.2, 6.5], y, theme)
        y += 0.34 * len(rows) + 0.35

    if scope.cost_plan.assumptions_register:
        af2 = _textbox(slide, theme.margin, y, theme.content_w, theme.slide_h - y - 0.7)
        _write(af2, "Assumptions register", size=13, bold=True, color=theme.blue_dark,
               first=True, space_after=5, theme=theme)
        for item in scope.cost_plan.assumptions_register:
            _write(af2, item, size=10, color=theme.muted, bullet=True, space_after=3, theme=theme)

    _footer(slide, footer, page, theme)


def _team_and_exclusions_slide(
    prs: Presentation, scope: ScopeOfWorkPayloadV2, deal_name: str | None,
    footer: str, page: int, theme: Theme = THEME,
) -> None:
    slide = _blank(prs)
    top = _slide_header(slide, "Team and exclusions", deal_name, theme)

    half = theme.content_w / 2 - 0.25
    team = scope.team_shape

    tf = _textbox(slide, theme.margin, top, half, 4.8)
    _write(tf, "Team", size=15, bold=True, color=theme.blue_dark, first=True,
           space_after=6, theme=theme)
    for member in team.core_team:
        _write(tf, member, size=10.5, color=theme.text, bullet=True, space_after=4, theme=theme)
    if team.specialists:
        _write(tf, "Specialists required", size=12, bold=True, color=theme.blue_dark,
               space_after=5, theme=theme)
        for s in team.specialists:
            _write(tf, s, size=10, color=theme.muted, bullet=True, space_after=3, theme=theme)

    xf = _textbox(slide, theme.margin + half + 0.5, top, half, 4.8)
    _write(xf, "Exclusions", size=15, bold=True, color=theme.blue_dark, first=True,
           space_after=4, theme=theme)
    _write(xf, "A scope that does not say what it excludes is not a scope.",
           size=9.5, color=theme.muted, italic=True, space_after=7, theme=theme)
    for exclusion in scope.exclusions:
        _write(xf, f"{exclusion.subject} — {exclusion.reason}", size=10,
               color=theme.text, bullet=True, space_after=4, theme=theme)

    _footer(slide, footer, page, theme)


# ------------------------------------------------------------------------- entry


def build_deck(
    scope: ScopeOfWorkPayloadV2,
    deal_name: str | None = None,
    version: int | None = None,
    theme: Theme = THEME,
) -> Presentation:
    """Assemble the deck. Split from `render_pptx` so tests can inspect the slides."""
    prs = open_presentation()
    prs.slide_width = Inches(theme.slide_w)
    prs.slide_height = Inches(theme.slide_h)

    footer = f"Library v{scope.library_version} · rules v{scope.rules_version}"
    if version is not None:
        footer = f"Version {version} · {footer}"

    _cover_slide(prs, scope, deal_name, version, theme)
    page = 2
    _engagement_slide(prs, scope, deal_name, footer, page, theme)
    page += 1
    page = _row_slides(prs, scope, deal_name, footer, page, theme)

    if scope.sequencing:
        _sequencing_slide(prs, scope, deal_name, footer, page, theme)
        page += 1

    _cost_slide(prs, scope, deal_name, footer, page, theme)
    page += 1
    _team_and_exclusions_slide(prs, scope, deal_name, footer, page, theme)

    return prs


def render_pptx(
    scope: ScopeOfWorkPayloadV2,
    deal_name: str | None = None,
    version: int | None = None,
) -> bytes:
    """The client-facing scope as a .pptx, returned as bytes."""
    buf = io.BytesIO()
    build_deck(scope, deal_name, version).save(buf)
    return buf.getvalue()


def render_pptx_from_payload(
    payload: dict[str, Any],
    deal_name: str | None = None,
    version: int | None = None,
) -> bytes:
    """Render a stored payload. Only schema v2 is exportable."""
    if payload.get("schema_version") != 2:
        raise ValueError("Only schema v2 scopes can be exported to PowerPoint")
    return render_pptx(ScopeOfWorkPayloadV2.model_validate(payload), deal_name, version)


def pptx_filename(deal_name: str, version: int) -> str:
    """`project-redline-scope-v2.pptx`, matching the PDF's naming."""
    slug = re.sub(r"[^a-z0-9]+", "-", deal_name.lower()).strip("-")
    return f"{slug}-scope-v{version}.pptx" if slug else f"scope-v{version}.pptx"
